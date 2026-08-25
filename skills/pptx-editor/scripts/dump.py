#!/usr/bin/env python3
"""PPTX 勘察：列出每页 shape 的名称、几何（英寸）、类型与文本/字体。

用法：
  python3 dump.py 文件.pptx                 # 全部页
  python3 dump.py 文件.pptx --slides 2,20   # 指定页
"""
import argparse

from pptx import Presentation
from pptx.util import Emu


def first_run(shape):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            return r
    return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--slides", help="逗号分隔页码，如 2,20；缺省全部")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    sel = None
    if args.slides:
        sel = {int(s) for s in args.slides.split(",") if s.strip()}

    print(f"slide size: {prs.slide_width} x {prs.slide_height} | slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        if sel and i not in sel:
            continue
        print(f"\n===== SLIDE {i} =====")
        for shape in slide.shapes:
            geo = ""
            try:
                geo = (f"({Emu(shape.left).inches:.2f},{Emu(shape.top).inches:.2f} "
                       f"{Emu(shape.width).inches:.2f}x{Emu(shape.height).inches:.2f})")
            except Exception:
                geo = "(no geo)"
            if shape.has_text_frame:
                text = shape.text_frame.text.replace("\n", "\\n")
                r = first_run(shape)
                f = ""
                if r is not None:
                    f = f" <{r.font.name or '?'}/{r.font.size.pt if r.font.size else '?'}pt"
                    if r.font.bold:
                        f += "/bold"
                    f += ">"
                print(f"  [{shape.name}] {geo} TEXT: {text[:120]}{f}")
            else:
                print(f"  [{shape.name}] {geo} ({shape.shape_type})")


if __name__ == "__main__":
    main()
