#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PPTX 精确修改 patch 脚本模板。

使用方法：复制本文件到 /tmp/fix_ppt_*.py，按任务填充 REPLACEMENTS / add_text 调用，
运行后逐条输出 OK/MISS，全部 OK 后再跑 verify.py 渲染验证。

要点：
- 按 shape.name 定位（先用 dump.py 勘察），不要按页内文字盲搜。
- run 级替换优先（保留格式）；一句话被拆成多个 run 时用 paragraph 级兜底。
- 新增元素照抄邻近元素的 font.name / size / bold / color / margins / line_spacing。
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PPT = "/path/to/演示文稿.pptx"

# 文本替换清单：(页码 1-based, 旧文本, 新文本)
REPLACEMENTS = [
    # (2, '旧表述', '新表述'),
    # (20, '旧来源', '新来源'),
]


def replace_in_paragraph(p, old, new):
    """paragraph 级兜底：拼接全部 run 文本，整体重写回第一个 run，保留其格式。"""
    full = "".join(r.text for r in p.runs)
    if old not in full:
        return False
    runs = p.runs
    runs[0].text = full.replace(old, new)
    for r in runs[1:]:
        r.text = ""
    return True


def rep_runs(slide, old, new, size_pt=None, wrap=None):
    """在指定 slide 内做 run 级替换（+ paragraph 级兜底）；返回命中数。"""
    n = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            if wrap is not None:
                shape.text_frame.word_wrap = wrap
            for p in shape.text_frame.paragraphs:
                done = False
                for r in p.runs:
                    if old in r.text:
                        r.text = r.text.replace(old, new)
                        if size_pt is not None:
                            r.font.size = Pt(size_pt)
                        n += 1
                        done = True
                if not done and old in "".join(r.text for r in p.runs):
                    if replace_in_paragraph(p, old, new):
                        n += 1
    return n


def add_text(slide, name, x, y, w, h, text, size, bold, hexcolor,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0,
             font_name="MiSans"):
    """新增文本框。font_name 等风格参数先用 dump.py 查同页邻近元素后照抄。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(hexcolor)
    return box


prs = Presentation(PPT)
report = []

for slide_idx, old, new in REPLACEMENTS:
    n = rep_runs(prs.slides[slide_idx - 1], old, new)
    report.append((slide_idx, old[:30], f"OK x{n}" if n else "MISS"))

# —— 新增元素示例（按需取消注释并修改）——
# s2 = prs.slides[1]
# add_text(s2, "p2_cap", 0.89, 3.60, 11.56, 0.30, "新增一行信号文字", 13, True, "17212B")
# report.append((2, "p2_cap 新增", "OK"))

# —— 几何/排版调整示例 ——
# for shape in prs.slides[19].shapes:
#     if shape.name == "p20_shift":
#         shape.top = Inches(4.71)
#         for p in shape.text_frame.paragraphs:
#             p.line_spacing = 1.0
# report.append((20, "p20_shift 调整", "OK"))

prs.save(PPT)

print("=== 修改结果 ===")
for slide, what, status in report:
    print(f"P{slide:>2} [{status}] {what}")
